import collections  # ensure collections.abc is loaded for python-pptx compat
import collections.abc  # pre-import to expose collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

TEMPLATE = 'CSG Overview.pptx'
OUTPUT = 'FlowForge-Sales-Deck.pptx'


def remove_all_slides(pres: Presentation):
    try:
        for idx in range(len(pres.slides) - 1, -1, -1):
            rId = pres.slides._sldIdLst[idx].rId
            pres.part.drop_rel(rId)
            pres.slides._sldIdLst.remove(pres.slides._sldIdLst[idx])
    except Exception:
        # if template is protected or structure differs, skip removal
        pass


def pick_content_layout(pres: Presentation):
    # prefer a layout that has title + content/body
    for layout in pres.slide_layouts:
        names = [ph.name.lower() for ph in layout.placeholders]
        has_title = any('title' in n for n in names)
        has_body = any(('content' in n) or ('body' in n) for n in names)
        if has_title and has_body:
            return layout
    # fallback: first layout
    return pres.slide_layouts[0]


def add_bullet_slide(pres: Presentation, layout, title: str, bullets: list[str], note: str | None):
    slide = pres.slides.add_slide(layout)
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title
    # Body
    # try to find a body/content placeholder (idx 1 typical)
    placeholder = None
    if len(slide.placeholders) > 1:
        placeholder = slide.placeholders[1]
    else:
        # fallback: add a text box
        left, top, width, height = Cm(2), Cm(5), Cm(24), Cm(10)
        placeholder = slide.shapes.add_textbox(left, top, width, height)
    tf = placeholder.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
    if note:
        slide.notes_slide.notes_text_frame.text = note
    return slide


def add_bsg_diagram(slide):
    left, top, width, height = Cm(2), Cm(10), Cm(24), Cm(3)
    w = width / 3 - Cm(0.5)
    labels = [
        ('Bronze', RGBColor(0x22, 0x88, 0xCC)),
        ('Silver', RGBColor(0x99, 0x99, 0x99)),
        ('Gold', RGBColor(0xDA, 0xA5, 0x20)),
    ]
    x = left
    for lbl, color in labels:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, w, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tx = shape.text_frame
        tx.text = lbl
        tx.paragraphs[0].font.bold = True
        tx.paragraphs[0].font.size = Pt(18)
        tx.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        x = x + w + Cm(0.75)
    # arrows
    slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + w + Cm(0.1), top + Cm(0.6), Cm(0.5), Cm(1.8))
    slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + 2 * w + Cm(1.35), top + Cm(0.6), Cm(0.5), Cm(1.8))


def add_three_pillars(slide):
    left, top, width, height = Cm(2), Cm(10), Cm(24), Cm(3)
    w = width / 3 - Cm(0.5)
    labels = ['AI', 'Medallion', 'Governance']
    x = left
    for lbl in labels:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, w, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        shape.line.color.rgb = RGBColor(0x88, 0x88, 0x88)
        tx = shape.text_frame
        tx.text = lbl
        tx.paragraphs[0].font.bold = True
        tx.paragraphs[0].font.size = Pt(18)
        x = x + w + Cm(0.75)


def add_2x2_grid(slide):
    left, top, width, height = Cm(4), Cm(9), Cm(18), Cm(8)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.background()
    rect.line.color.rgb = RGBColor(0x88, 0x88, 0x88)
    # mid-lines
    slide.shapes.add_line(left + width / 2, top, left + width / 2, top + height)
    slide.shapes.add_line(left, top + height / 2, left + width, top + height / 2)
    # labels
    tx1 = slide.shapes.add_textbox(left + width - Cm(6), top - Cm(1.2), Cm(6), Cm(1))
    tx1.text_frame.text = 'Ease of Use →'
    tx2 = slide.shapes.add_textbox(left - Cm(3.5), top + Cm(1), Cm(3.5), Cm(1))
    tx2.text_frame.text = 'Control / Neutrality ↑'
    # marker upper-right
    marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + width * 0.70, top + height * 0.10, Cm(2.0), Cm(1.2))
    marker.fill.solid(); marker.fill.fore_color.rgb = RGBColor(0x22, 0x88, 0xCC)
    marker.line.color.rgb = RGBColor(0x22, 0x88, 0xCC)
    marker.text_frame.text = 'FlowForge'
    marker.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_hub_spoke(slide):
    center_x, center_y = Cm(15), Cm(11)
    hub = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, center_x - Cm(2.5), center_y - Cm(1), Cm(5), Cm(2))
    hub.text_frame.text = 'FlowForge'
    hub.fill.solid(); hub.fill.fore_color.rgb = RGBColor(0x22, 0x88, 0xCC)
    hub.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    spokes = ['S3', 'ADLS', 'Snowflake', 'Databricks', 'Redshift', 'BigQuery', 'Power BI', 'Tableau', 'Postgres', 'SQL Server', 'MySQL']
    import math
    r = Cm(7)
    for i, label in enumerate(spokes):
        ang = 2 * math.pi * i / len(spokes)
        x = center_x + r * math.cos(ang)
        y = center_y + r * math.sin(ang)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - Cm(2.2), y - Cm(0.8), Cm(4.4), Cm(1.6))
        box.text_frame.text = label
        slide.shapes.add_connector(1, center_x, center_y, x, y)


def main():
    pres = Presentation(TEMPLATE)
    remove_all_slides(pres)
    layout = pick_content_layout(pres)

    slides_def = [
        ('FlowForge — Modern Data Orchestration Platform', [
            'Self‑service data pipelines that run in your cloud',
            'Build production‑grade pipelines without heavy engineering',
            'AI‑assisted schema detection and configuration',
            'Vendor‑neutral; works with your existing stack',
            'Simple to roll out and operate under your controls'
        ], 'Think “Power BI for data pipelines,” with cloud‑native control in your account.'),
        ('Why Data Plumbing Slows Business', [
            'Tool sprawl and engineering bottlenecks stall delivery',
            'DIY code is hard to govern, maintain, and scale',
            'Analysts wait; leaders question data consistency',
            'Compliance and data residency add friction'
        ], 'Ask which pain resonates most across their teams.'),
        ('The Case for a Simpler, Open Approach', [
            'Cloud‑hosted yet cloud‑agnostic (AWS/Azure/GCP)',
            'Standard patterns (Medallion) with built‑in governance',
            'Lighter operations; lower cognitive and cost overhead',
            'Fit with current tools; no rip‑and‑replace'
        ], 'Simplicity is the shortest path to impact.'),
        ('What Is FlowForge (In Your Cloud)', [
            'Build, run, and govern pipelines without code',
            'AI‑powered configuration and schema detection',
            'Environment/team isolation and RBAC',
            'Built‑in metadata catalog and monitoring'
        ], 'A focused platform that fits your operating model.'),
        ('How It Works — Bronze → Silver → Gold', [
            'Bronze: land and standardize files with audit columns',
            'Silver: cleanse, deduplicate, apply mappings/PKs',
            'Gold: publish compressed, analytics‑ready datasets; catalog',
            'Triggers: manual, scheduled, and dependency‑based'
        ], 'Opinionated best practices, end‑to‑end.'),
        ('Six Pillars of FlowForge', [
            'AI‑assisted setup: schemas, column names, PK suggestions',
            'File processing at scale: pattern matching and multi‑file automation',
            'Orchestration: multi‑env, team isolation, dependency triggers',
            'Monitoring: live status, logs, metrics per workflow/job',
            'Catalog: automatic metadata, schema, lineage',
            'Vendor‑neutral I/O: S3/ADLS, Snowflake, Databricks, Redshift, BigQuery'
        ], 'Everything needed without platform bloat.'),
        ('Why We Win', [
            'Self‑service for power users; engineers focus on complex work',
            'Runs in your cloud; data stays in‑account for compliance',
            'Neutral by design; keep choice of lake/warehouse',
            'Predictable commercial model; avoid per‑row surprises'
        ], 'The sweet spot between flexibility and ease.'),
        ('Security & Deployment (Your Cloud)', [
            'Data residency: artifacts live in your S3/ADLS',
            'Identity: IAM/Key Vault integration; role‑based access and approvals',
            'Encryption & networking: KMS/CMK, private networking, logging, audits',
            'Works with your monitoring and ticketing tools'
        ], 'Keep your controls; FlowForge fits into them.'),
        ('What Good Looks Like (Outcomes)', [
            '50–80% less effort for common pipelines',
            '30–60% of routine requests shift to self‑service',
            'Consolidate tools; meaningfully reduce annual spend',
            'Consistent patterns increase data trust and reuse'
        ], 'Tailor with your metrics; direction is consistent.'),
        ('Works With Your Stack (Ecosystem)', [
            'Storage: Amazon S3, Azure Blob/ADLS',
            'Warehouses/Lakes: Snowflake, Databricks, Redshift, BigQuery',
            'Databases: PostgreSQL, SQL Server, MySQL (extensible)',
            'Analytics: Power BI, Tableau, Looker via standards‑based outputs'
        ], 'We complement your stack, not replace it.'),
        ('Predictable and Simple (Packaging)', [
            'Predictable tiers; no per‑row billing complexity',
            'Runs in your cloud — you control infra costs',
            'Licensing aligned to environments/teams, not seats',
            'Optional low‑lift pilot to validate fit'
        ], 'Emphasize clarity and control.'),
        ('What’s Next (Highlights; No Dates)', [
            'Additional database and API connectors',
            'Data quality rules and monitoring',
            'Document/semi‑structured data processing',
            'Guardrails, approvals, and policy automation',
            'Optional streaming integrations as needed'
        ], 'Evolution without lock‑in or forced rewrites.'),
        ('See It on Your Data (Low‑Lift Pilot)', [
            'Bring one or two representative pipelines',
            'Run in your cloud with your security controls',
            'Keep outputs in your storage for immediate value',
            'Clear success criteria; minimal overhead'
        ], 'Prove value simply, then expand.'),
    ]

    added = []
    for title, bullets, note in slides_def:
        added.append(add_bullet_slide(pres, layout, title, bullets, note))

    # add diagrams to selected slides
    try:
        add_hub_spoke(added[0])        # Slide 1
        add_three_pillars(added[3])    # Slide 4
        add_bsg_diagram(added[4])      # Slide 5
        add_2x2_grid(added[6])         # Slide 7
    except Exception:
        pass

    pres.save(OUTPUT)
    print(f'Generated {OUTPUT}')


if __name__ == '__main__':
    main()
